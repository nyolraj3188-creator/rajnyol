"""
PDF Tools backend service.
Provides server-side conversions and security operations that cannot run
purely in the browser. Each endpoint accepts an uploaded file (multipart)
and streams back the processed result.
"""
import os
import shutil
import uuid
import subprocess
import tempfile
import difflib
from pathlib import Path

from fastapi import APIRouter, UploadFile, File, Form, HTTPException, BackgroundTasks
from fastapi.responses import FileResponse, JSONResponse

router = APIRouter(prefix="/api/pdf", tags=["pdf-tools"])

WORK_ROOT = Path(tempfile.gettempdir()) / "pdfpro_jobs"
WORK_ROOT.mkdir(parents=True, exist_ok=True)

MIME = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def _new_job():
    d = WORK_ROOT / uuid.uuid4().hex
    d.mkdir(parents=True, exist_ok=True)
    return d


def _cleanup(path: Path):
    shutil.rmtree(path, ignore_errors=True)


async def _save_upload(upload: UploadFile, job: Path) -> Path:
    dest = job / (upload.filename or "input")
    with open(dest, "wb") as f:
        while chunk := await upload.read(1024 * 1024):
            f.write(chunk)
    await upload.seek(0)
    return dest


def _libreoffice(src: Path, target: str, outdir: Path) -> Path:
    """Convert a file to `target` (e.g. 'pdf', 'docx') via headless LibreOffice."""
    profile = outdir / "lo_profile"
    cmd = [
        "soffice", "--headless", "--norestore", "--nolockcheck",
        f"-env:UserInstallation=file://{profile}",
        "--convert-to", target, "--outdir", str(outdir), str(src),
    ]
    proc = subprocess.run(cmd, capture_output=True, timeout=180)
    if proc.returncode != 0:
        raise HTTPException(500, f"Conversion failed: {proc.stderr.decode()[:300]}")
    ext = target.split(":")[0]
    out = outdir / (src.stem + "." + ext)
    if not out.exists():
        cands = list(outdir.glob(f"{src.stem}.*"))
        cands = [c for c in cands if c.suffix.lstrip(".") == ext]
        if not cands:
            raise HTTPException(500, "Converted file not produced.")
        out = cands[0]
    return out


def _respond(path: Path, filename: str, ext: str, job: Path, bg: BackgroundTasks):
    bg.add_task(_cleanup, job)
    return FileResponse(str(path), media_type=MIME.get(ext, "application/octet-stream"), filename=filename, background=bg)


@router.get("/health")
async def health():
    tools = {}
    for t in ["soffice", "gs", "qpdf", "tesseract", "pdftoppm"]:
        tools[t] = shutil.which(t) is not None
    try:
        import ocrmypdf  # noqa
        tools["ocrmypdf"] = True
    except Exception:
        tools["ocrmypdf"] = False
    return {"ok": True, "tools": tools}


# ---------- Office / HTML -> PDF ----------
@router.post("/office-to-pdf")
async def office_to_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        out = _libreoffice(src, "pdf", job)
        name = Path(file.filename).stem + ".pdf"
        return _respond(out, name, "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, str(e))


@router.post("/html-to-pdf")
async def html_to_pdf(background_tasks: BackgroundTasks, url: str = Form(None), html: str = Form(None)):
    import requests
    job = _new_job()
    try:
        content = html
        if url:
            r = requests.get(url, timeout=30, headers={"User-Agent": "Mozilla/5.0"})
            r.raise_for_status()
            content = r.text
        if not content:
            raise HTTPException(400, "Provide a URL or HTML content.")
        src = job / "page.html"
        src.write_text(content, encoding="utf-8")
        out = _libreoffice(src, "pdf", job)
        return _respond(out, "webpage.pdf", "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, str(e))


def _has_text_layer(src: Path) -> bool:
    """True if the PDF contains a meaningful extractable text layer."""
    import pdfplumber
    try:
        with pdfplumber.open(str(src)) as pdf:
            for page in pdf.pages[:5]:
                if len((page.extract_text() or "").strip()) > 20:
                    return True
    except Exception:
        pass
    return False


def _ocr_if_scanned(src: Path, job: Path, lang: str = "eng") -> Path:
    """If the PDF is image-only (scanned), OCR it so downstream converters see text."""
    if _has_text_layer(src):
        return src
    _require_ocr_tools()
    import ocrmypdf
    ocred = job / (src.stem + "_ocred.pdf")
    try:
        ocrmypdf.ocr(str(src), str(ocred), skip_text=True, language=lang, optimize=0, progress_bar=False)
        if ocred.exists():
            return ocred
    except Exception:
        pass
    return src


def _scanned_pdf_to_docx(src: Path, job: Path, out: Path, lang: str):
    """OCR a scanned PDF and build an editable Word document from the recognized text."""
    import re
    _require_ocr_tools()
    import ocrmypdf
    from docx import Document
    sidecar = job / "ocr_text.txt"
    ocred = job / (src.stem + "_ocred.pdf")
    ocrmypdf.ocr(str(src), str(ocred), force_ocr=True, language=lang,
                 sidecar=str(sidecar), optimize=0, progress_bar=False)
    text = sidecar.read_text(encoding="utf-8", errors="ignore") if sidecar.exists() else ""
    if not text.strip():
        raise HTTPException(422, "No readable text found in this scanned PDF.")
    illegal = re.compile(r"[\x00-\x08\x0b\x0e-\x1f\x7f-\x9f\ufffe\uffff]")
    doc = Document()
    pages = text.split("\f")
    for pi, page in enumerate(pages):
        for line in page.splitlines():
            clean = illegal.sub("", line).strip()
            if clean:
                doc.add_paragraph(clean)
        if pi < len(pages) - 1 and page.strip():
            doc.add_page_break()
    doc.save(str(out))


def _require_ocr_tools():
    """Ensure the OCR toolchain (Tesseract + Ghostscript) is available on PATH."""
    missing = [t for t in ("tesseract", "gs") if shutil.which(t) is None]
    if missing:
        raise HTTPException(
            503,
            "OCR is temporarily unavailable on the server (missing: "
            + ", ".join(missing) + "). Please try again shortly.",
        )


def _zip_has_parts(path: Path, required: list) -> bool:
    """A valid OOXML file is a zip that contains the required internal parts
    and is not corrupt."""
    import zipfile
    try:
        if not path.exists() or path.stat().st_size < 200:
            return False
        if not zipfile.is_zipfile(str(path)):
            return False
        with zipfile.ZipFile(str(path)) as z:
            names = set(z.namelist())
            if not set(required).issubset(names):
                return False
            if z.testzip() is not None:  # first corrupt member, if any
                return False
        return True
    except Exception:
        return False


def _docx_is_valid(path: Path) -> bool:
    """Strict validity check so we never hand back a file Word refuses to open."""
    if not _zip_has_parts(path, ["[Content_Types].xml", "word/document.xml"]):
        return False
    try:
        from docx import Document
        doc = Document(str(path))
        # touch the body so a structurally-broken document raises here
        _ = doc.paragraphs
        return True
    except Exception:
        return False


def _xlsx_is_valid(path: Path) -> bool:
    """Strict validity check for Excel output."""
    if not _zip_has_parts(path, ["[Content_Types].xml", "xl/workbook.xml"]):
        return False
    try:
        from openpyxl import load_workbook
        wb = load_workbook(str(path), read_only=True)
        _ = wb.sheetnames
        wb.close()
        return True
    except Exception:
        return False


def _normalize_office(src: Path, ext: str, job: Path) -> Path:
    """Re-save an Office file through headless LibreOffice so it opens cleanly
    in Microsoft Word / Excel. Returns the normalized file, or the original if
    LibreOffice is unavailable or the round-trip fails."""
    if shutil.which("soffice") is None:
        return src
    try:
        norm_dir = job / f"norm_{ext}"
        norm_dir.mkdir(exist_ok=True)
        produced = _libreoffice(src, ext, norm_dir)
        final = job / (src.stem + "_ms." + ext)
        shutil.copy(str(produced), str(final))
        return final
    except Exception:
        return src


# ---------- PDF -> Word ----------
@router.post("/pdf-to-word")
async def pdf_to_word(background_tasks: BackgroundTasks, file: UploadFile = File(...), lang: str = Form("eng")):
    from pdf2docx import Converter
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        out = job / (src.stem + ".docx")
        if _has_text_layer(src):
            try:
                cv = Converter(str(src))
                cv.convert(str(out))
                cv.close()
            except Exception:
                out.unlink(missing_ok=True)
            if not out.exists() or not _docx_is_valid(out):
                out.unlink(missing_ok=True)
                _scanned_pdf_to_docx(src, job, out, lang)
        else:
            _scanned_pdf_to_docx(src, job, out, lang)
        if not out.exists() or not _docx_is_valid(out):
            raise HTTPException(500, "Conversion produced an unreadable document. Please try another file.")
        # Normalize through LibreOffice so the .docx opens cleanly in MS Word.
        final = _normalize_office(out, "docx", job)
        if _docx_is_valid(final):
            out = final
        return _respond(out, src.stem + ".docx", "docx", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"PDF to Word failed: {e}")


# ---------- PDF -> Excel ----------
@router.post("/pdf-to-excel")
async def pdf_to_excel(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    import pdfplumber
    from openpyxl import Workbook
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        src = _ocr_if_scanned(src, job)
        wb = Workbook()
        wb.remove(wb.active)
        found = False
        with pdfplumber.open(str(src)) as pdf:
            for i, page in enumerate(pdf.pages):
                ws = wb.create_sheet(title=f"Page {i+1}"[:31])
                tables = page.extract_tables()
                if tables:
                    found = True
                    for tbl in tables:
                        for row in tbl:
                            ws.append([("" if c is None else c) for c in row])
                        ws.append([])
                else:
                    text = page.extract_text() or ""
                    for line in text.splitlines():
                        ws.append([line])
        if not wb.sheetnames:
            wb.create_sheet("Sheet1")
        out = job / (src.stem + ".xlsx")
        wb.save(str(out))
        if not _xlsx_is_valid(out):
            raise HTTPException(500, "Conversion produced an unreadable spreadsheet. Please try another file.")
        # Normalize through LibreOffice so the .xlsx opens cleanly in MS Excel.
        final = _normalize_office(out, "xlsx", job)
        if _xlsx_is_valid(final):
            out = final
        return _respond(out, src.stem + ".xlsx", "xlsx", job, background_tasks)
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"PDF to Excel failed: {e}")


# ---------- PDF -> PowerPoint (one image per slide) ----------
@router.post("/pdf-to-ppt")
async def pdf_to_ppt(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    from pdf2image import convert_from_path
    from pptx import Presentation
    from pptx.util import Inches
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        images = convert_from_path(str(src), dpi=120)
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]
        for idx, img in enumerate(images):
            p = job / f"s{idx}.png"
            img.save(str(p), "PNG")
            slide = prs.slides.add_slide(blank)
            iw, ih = img.size
            ratio = min(prs.slide_width / iw, prs.slide_height / ih)
            w = int(iw * ratio); h = int(ih * ratio)
            left = int((prs.slide_width - w) / 2); top = int((prs.slide_height - h) / 2)
            slide.shapes.add_picture(str(p), left, top, width=w, height=h)
        out = job / (src.stem + ".pptx")
        prs.save(str(out))
        return _respond(out, src.stem + ".pptx", "pptx", job, background_tasks)
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"PDF to PowerPoint failed: {e}")


# ---------- OCR ----------
@router.post("/ocr")
async def ocr_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...), lang: str = Form("eng")):
    import ocrmypdf
    job = _new_job()
    try:
        _require_ocr_tools()
        src = await _save_upload(file, job)
        out = job / (src.stem + "_ocr.pdf")
        try:
            ocrmypdf.ocr(str(src), str(out), skip_text=True, language=lang, optimize=0, progress_bar=False)
        except Exception as oe:
            raise HTTPException(500, f"OCR failed: {str(oe)[:300]}")
        if not out.exists():
            raise HTTPException(500, "OCR produced no output.")
        return _respond(out, src.stem + "_ocr.pdf", "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, str(e))


# ---------- Protect (encrypt) ----------
@router.post("/protect")
async def protect_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...), password: str = Form(...)):
    import pikepdf
    job = _new_job()
    try:
        if not password:
            raise HTTPException(400, "Password required.")
        src = await _save_upload(file, job)
        out = job / (src.stem + "_protected.pdf")
        with pikepdf.open(str(src)) as pdf:
            pdf.save(str(out), encryption=pikepdf.Encryption(owner=password, user=password, R=6))
        return _respond(out, src.stem + "_protected.pdf", "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"Protect failed: {e}")


# ---------- Unlock (decrypt) ----------
@router.post("/unlock")
async def unlock_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...), password: str = Form("")):
    import pikepdf
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        out = job / (src.stem + "_unlocked.pdf")
        try:
            with pikepdf.open(str(src), password=password) as pdf:
                pdf.save(str(out))
        except pikepdf.PasswordError:
            raise HTTPException(400, "Wrong password for this PDF.")
        return _respond(out, src.stem + "_unlocked.pdf", "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"Unlock failed: {e}")


# ---------- Repair (ghostscript rewrite) ----------
@router.post("/repair")
async def repair_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        out = job / (src.stem + "_repaired.pdf")
        cmd = ["gs", "-o", str(out), "-sDEVICE=pdfwrite", "-dPDFSTOPONERROR=false", str(src)]
        proc = subprocess.run(cmd, capture_output=True, timeout=180)
        if not out.exists():
            raise HTTPException(500, f"Repair failed: {proc.stderr.decode()[:300]}")
        return _respond(out, src.stem + "_repaired.pdf", "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, str(e))


# ---------- PDF/A ----------
@router.post("/pdfa")
async def pdf_to_pdfa(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        out = job / (src.stem + "_pdfa.pdf")
        cmd = ["gs", "-dPDFA=2", "-dBATCH", "-dNOPAUSE", "-sColorConversionStrategy=UseDeviceIndependentColor",
               "-sDEVICE=pdfwrite", "-dPDFACompatibilityPolicy=1", f"-sOutputFile={out}", str(src)]
        proc = subprocess.run(cmd, capture_output=True, timeout=180)
        if not out.exists():
            raise HTTPException(500, f"PDF/A conversion failed: {proc.stderr.decode()[:300]}")
        return _respond(out, src.stem + "_pdfa.pdf", "pdf", job, background_tasks)
    except HTTPException:
        _cleanup(job); raise
    except Exception as e:
        _cleanup(job); raise HTTPException(500, str(e))


# ---------- Crop (trim margins by percent) ----------
@router.post("/crop")
async def crop_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...), margin: float = Form(5.0)):
    from pypdf import PdfReader, PdfWriter
    job = _new_job()
    try:
        src = await _save_upload(file, job)
        reader = PdfReader(str(src))
        writer = PdfWriter()
        m = max(0.0, min(40.0, margin)) / 100.0
        for page in reader.pages:
            box = page.mediabox
            w = float(box.width); h = float(box.height)
            page.cropbox.lower_left = (float(box.left) + w * m, float(box.bottom) + h * m)
            page.cropbox.upper_right = (float(box.right) - w * m, float(box.top) - h * m)
            writer.add_page(page)
        out = job / (src.stem + "_cropped.pdf")
        with open(out, "wb") as f:
            writer.write(f)
        return _respond(out, src.stem + "_cropped.pdf", "pdf", job, background_tasks)
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"Crop failed: {e}")


# ---------- Compare two PDFs (text diff) ----------
@router.post("/compare")
async def compare_pdf(background_tasks: BackgroundTasks, file1: UploadFile = File(...), file2: UploadFile = File(...)):
    import pdfplumber
    job = _new_job()
    try:
        p1 = await _save_upload(file1, job)
        p2 = await _save_upload(file2, job)

        def text_of(p):
            lines = []
            with pdfplumber.open(str(p)) as pdf:
                for page in pdf.pages:
                    lines.extend((page.extract_text() or "").splitlines())
            return lines

        a, b = text_of(p1), text_of(p2)
        sm = difflib.SequenceMatcher(None, a, b)
        rows = []
        for tag, i1, i2, j1, j2 in sm.get_opcodes():
            if tag == "equal":
                for line in a[i1:i2]:
                    rows.append({"type": "equal", "text": line})
            elif tag == "replace":
                for line in a[i1:i2]:
                    rows.append({"type": "removed", "text": line})
                for line in b[j1:j2]:
                    rows.append({"type": "added", "text": line})
            elif tag == "delete":
                for line in a[i1:i2]:
                    rows.append({"type": "removed", "text": line})
            elif tag == "insert":
                for line in b[j1:j2]:
                    rows.append({"type": "added", "text": line})
        ratio = round(sm.ratio() * 100, 1)
        _cleanup(job)
        return JSONResponse({"similarity": ratio, "rows": rows[:1000]})
    except Exception as e:
        _cleanup(job); raise HTTPException(500, f"Compare failed: {e}")
